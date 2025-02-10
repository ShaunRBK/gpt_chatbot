import os
import uuid
import json
from io import BytesIO
import base64

from openai import OpenAI
from flask import Flask, render_template, request, jsonify, send_file

# Use a non-interactive backend for Matplotlib
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# Additional imports for file processing and document generation
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
from fpdf import FPDF

# For RAG (retrieval-augmented generation)
import faiss
import numpy as np

# Initialize the OpenAI client (using your migrated code)
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Initialize Flask
app = Flask(__name__)

# -----------------------------
# Global Variables
# -----------------------------
faiss_index = None          # FAISS index for vector search (from uploaded documents)
chunk_texts = []            # List mapping each vector in FAISS to its corresponding text chunk
full_document_text = None   # Full text from the uploaded document (for fallback)
uploaded_df = None          # For data analysis, store the uploaded DataFrame here.
# (In production, consider using session or database storage rather than globals.)

# -----------------------------
# Helper Functions for File Extraction and RAG
# -----------------------------

def extract_text_from_pdf(file):
    """
    Extracts and returns text from a PDF file.

    Args:
      file (FileStorage): The uploaded PDF file.

    Returns:
      str: The extracted text.
    """
    file.seek(0)
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_text_from_docx(file):
    """
    Extracts and returns text from a DOCX (Word) file.

    Args:
      file (FileStorage): The uploaded DOCX file.

    Returns:
      str: The extracted text.
    """
    file.seek(0)
    doc = docx.Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

def extract_text_from_excel(file, sheet=None):
    """
    Extracts text from an Excel file by converting its content to a string.
    Also reads the file into a DataFrame for data analysis.

    Args:
      file (FileStorage): The uploaded Excel file.
      sheet (str, optional): The name of the sheet to use. If None, defaults to the first sheet.

    Returns:
      str: The content of the Excel file as text.
    """
    file.seek(0)
    try:
        # If sheet is provided, read that sheet; otherwise, read the first sheet.
        df = pd.read_excel(file, sheet_name=sheet) if sheet else pd.read_excel(file)
        global uploaded_df
        uploaded_df = df.copy()  # Save for later analysis
        return df.to_string()
    except Exception as e:
        return f"Error reading Excel file: {e}"

def extract_text_from_csv(file):
    """
    Extracts text from a CSV file by reading it into a DataFrame.

    Args:
      file (FileStorage): The uploaded CSV file.

    Returns:
      str: The CSV content as text.
    """
    file.seek(0)
    try:
        df = pd.read_csv(file)
        global uploaded_df
        uploaded_df = df.copy()
        return df.to_string()
    except Exception as e:
        return f"Error reading CSV file: {e}"

def extract_text_from_pptx(file):
    """
    Extracts and returns text from a PowerPoint (.pptx) file.

    Args:
      file (FileStorage): The uploaded PowerPoint file.

    Returns:
      str: The extracted text.
    """
    file.seek(0)
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def chunk_text(text, chunk_size=500):
    """
    Splits the provided text into smaller chunks based on a target number of words.

    Args:
      text (str): The full text to be chunked.
      chunk_size (int): Maximum words per chunk (default is 500).

    Returns:
      list[str]: A list of text chunks.
    """
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i+chunk_size])
        chunks.append(chunk)
    return chunks

def create_vector_store(text):
    """
    Chunks the provided text, computes embeddings for each chunk using OpenAI's embedding model,
    and builds a FAISS vector index for retrieval.

    Also stores the full text for fallback.

    Args:
      text (str): The full document text.

    Effects:
      Updates globals: 'faiss_index', 'chunk_texts', and 'full_document_text'.
    """
    global faiss_index, chunk_texts, full_document_text
    chunks = chunk_text(text)
    d = 1536  # Embedding dimension for text-embedding-ada-002
    index = faiss.IndexFlatL2(d)
    embeddings = []
    chunk_texts = []  # Reset the list of chunks
    for chunk in chunks:
        try:
            response = client.embeddings.create(model="text-embedding-ada-002", input=chunk)
            embedding = np.array(response.data[0].embedding, dtype='float32')
            embeddings.append(embedding)
            chunk_texts.append(chunk)
        except Exception as e:
            print(f"Error computing embedding for a chunk: {e}")
    if embeddings:
        embeddings_np = np.vstack(embeddings)
        index.add(embeddings_np)
        faiss_index = index
    full_document_text = text

def retrieve_relevant_chunks(query, top_k=5):
    """
    Computes the embedding for the query and retrieves the top_k most similar text chunks from the FAISS index.

    Args:
      query (str): The user query.
      top_k (int): Number of top similar chunks to retrieve (default is 5).

    Returns:
      list[str]: Relevant text chunks. If none found, returns the full document text.
    """
    global faiss_index, chunk_texts, full_document_text
    if faiss_index is None:
        return []
    try:
        response = client.embeddings.create(model="text-embedding-ada-002", input=query)
        query_embedding = np.array(response.data[0].embedding, dtype='float32').reshape(1, -1)
        distances, indices = faiss_index.search(query_embedding, top_k)
        relevant_chunks = []
        for idx in indices[0]:
            if idx < len(chunk_texts):
                relevant_chunks.append(chunk_texts[idx])
        if not relevant_chunks and full_document_text:
            relevant_chunks = [full_document_text]
        return relevant_chunks
    except Exception as e:
        print(f"Error retrieving relevant chunks: {e}")
        return []

# -----------------------------
# Helper Functions for Document Generation (unchanged)
# -----------------------------

def generate_docx(text):
    """
    Generates a Microsoft Word (.docx) document containing the provided text.

    Args:
      text (str): The content to be included in the document.

    Returns:
      BytesIO: A stream containing the generated DOCX file.
    """
    document = docx.Document()
    document.add_heading("Generated Document", level=1)
    for paragraph in text.split("\n"):
        document.add_paragraph(paragraph)
    output = BytesIO()
    document.save(output)
    output.seek(0)
    return output

def generate_pptx(text):
    """
    Generates a PowerPoint (.pptx) presentation containing the provided text.

    Args:
      text (str): The content to be included in the presentation.

    Returns:
      BytesIO: A stream containing the generated PPTX file.
    """
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Generated Presentation"
    slide.placeholders[1].text = text
    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output

def generate_pdf(text):
    """
    Generates a PDF document containing the provided text.

    Args:
      text (str): The content to be included in the PDF.

    Returns:
      BytesIO: A stream containing the generated PDF file.
    """
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in text.split("\n"):
        pdf.multi_cell(0, 10, line)
    pdf_bytes = pdf.output(dest="S").encode("latin1")
    output = BytesIO(pdf_bytes)
    output.seek(0)
    return output

def store_file_temporarily(file_stream, ext):
    """
    Stores a file stream in temporary storage and returns a unique download URL.

    Args:
      file_stream (BytesIO): The file stream.
      ext (str): The file extension (e.g., 'docx', 'pptx', 'pdf', or 'png').

    Returns:
      str: A relative download URL for the stored file.
    """
    os.makedirs("temp_files", exist_ok=True)
    temp_filename = f"{uuid.uuid4()}.{ext}"
    temp_path = os.path.join("temp_files", temp_filename)
    with open(temp_path, "wb") as f:
        f.write(file_stream.read())
    return f"/download/{temp_filename}"

# -----------------------------
# Helper Function for Data Analysis (updated)
# -----------------------------

def analyze_data(query):
    """
    Analyzes the uploaded data (stored in the global 'uploaded_df') based on the query.

    If the query mentions "plot" (case-insensitive), it will attempt to generate a plot.
      - If there are numeric columns, it generates a histogram.
      - Otherwise, if it finds candidate categorical columns (e.g. exercise types),
        it creates a bar plot of value counts.
    Otherwise, it returns summary statistics.

    Args:
      query (str): The analysis query regarding the uploaded data.

    Returns:
      dict: A dictionary with key "analysis_result" (string) and optionally "plot_inline" (string, an HTML img tag with the plot).
    """
    global uploaded_df
    result = {}
    if uploaded_df is None:
        result["analysis_result"] = "No data available for analysis."
        return result

    if "plot" in query.lower():
        # Try numeric columns first.
        numeric_cols = uploaded_df.select_dtypes(include=["number"]).columns.tolist()
        plt.figure(figsize=(8, 6))
        if numeric_cols:
            # Plot histograms for each numeric column on the same figure.
            for col in numeric_cols:
                plt.hist(uploaded_df[col].dropna(), bins=20, alpha=0.5, label=str(col))
            plt.legend()
            plt.title("Histogram of Numeric Columns")
        else:
            # If no numeric columns, try categorical columns.
            cat_cols = [col for col in uploaded_df.select_dtypes(include=["object"]).columns
                        if uploaded_df[col].nunique() > 1]
            if not cat_cols:
                result["analysis_result"] = "No appropriate data available for plotting."
                return result
            # Choose the first categorical column.
            col = cat_cols[0]
            counts = uploaded_df[col].value_counts()
            counts.plot(kind="bar")
            plt.title(f"Bar Plot of {col}")
            plt.xlabel(col)
            plt.ylabel("Count")
        plt.tight_layout()
        # Save the plot to a BytesIO stream.
        plot_stream = BytesIO()
        plt.savefig(plot_stream, format="png")
        plt.close()
        plot_stream.seek(0)
        # Encode the image in base64.
        plot_bytes = plot_stream.getvalue()
        encoded_plot = base64.b64encode(plot_bytes).decode('utf-8')
        img_tag = f'<img src="data:image/png;base64,{encoded_plot}" style="max-width: 100%;">'
        result["analysis_result"] = "Plot generated based on the uploaded data."
        result["plot_inline"] = img_tag
    else:
        # Return summary statistics including non-numeric columns.
        result["analysis_result"] = uploaded_df.describe(include="all").to_string()
    return result
# -----------------------------
# Flask Endpoints
# -----------------------------

@app.route("/")
def index():
    """
    Renders the main UI page.
    """
    return render_template("index.html")

@app.route("/chat", methods=["POST"])
def chat():
    """
    Handles user messages from the main prompt.

    Depending on the user's intent, the assistant might call one of the following functions:
      - generate_document: To generate a document file.
      - analyze_data: To analyze uploaded data.

    Returns:
      JSON: A response with either a regular chat answer or a download link/analysis result.
    """
    data = request.get_json()
    user_message = data.get("message")
    if not user_message:
        return jsonify({"error": "No message provided"}), 400

    try:
        # Build conversation context.
        messages = [
            {"role": "system", "content": (
                "You are a helpful assistant. "
                "If the user's request indicates a need to generate a document "
                "(for example, a Word document, PowerPoint, or PDF), respond by calling the function "
                "generate_document with a JSON object containing 'file_type' and 'document_content'. "
                "If the user's request is about data analysis of an uploaded Excel/CSV file, call the function "
                "analyze_data with a JSON object containing 'query'. "
                "For Excel files with multiple sheets, the user can include the sheet name in their query "
                "(for example, 'use sheet SalesData'). "
                "Otherwise, answer normally."
            )}
        ]
        # (Optional) Add retrieval-based context if available.
        relevant_chunks = retrieve_relevant_chunks(user_message, top_k=5)
        context_text = "\n\n".join(relevant_chunks)
        if context_text:
            messages.append({"role": "system", "content": f"Relevant document context:\n{context_text}"})
        messages.append({"role": "user", "content": user_message})

        # Prepare function definitions.
        functions = []
        functions.append({
            "name": "generate_document",
            "description": (
                "Generates a document. The output should be a JSON with two keys: "
                "'file_type' (one of ['docx', 'pptx', 'pdf']) and "
                "'document_content' (the text content for the document)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_type": {"type": "string", "enum": ["docx", "pptx", "pdf"]},
                    "document_content": {"type": "string", "description": "The text content for the document."}
                },
                "required": ["file_type", "document_content"]
            }
        })
        if uploaded_df is not None:
            functions.append({
                "name": "analyze_data",
                "description": (
                    "Analyzes the uploaded data based on the user's query. "
                    "If the query mentions 'plot', a plot image is generated; otherwise, summary statistics are returned. "
                    "The output should be a JSON with key 'analysis_result' (string) and optionally 'plot_url' (string)."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string", "description": "The analysis query regarding the uploaded data."}
                    },
                    "required": ["query"]
                }
            })

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            functions=functions,
            function_call="auto"
        )
        message_response = response.choices[0].message
        if hasattr(message_response, "to_dict"):
            message_response = message_response.to_dict()

        if message_response.get("function_call"):
            function_call = message_response["function_call"]
            fname = function_call.get("name")
            args = json.loads(function_call.get("arguments", "{}"))
            if fname == "generate_document":
                file_type = args.get("file_type")
                document_content = args.get("document_content")
                if not file_type or not document_content:
                    return jsonify({"answer": "Failed to parse document generation parameters."})
                if file_type == "docx":
                    file_stream = generate_docx(document_content)
                    ext = "docx"
                elif file_type == "pptx":
                    file_stream = generate_pptx(document_content)
                    ext = "pptx"
                elif file_type == "pdf":
                    file_stream = generate_pdf(document_content)
                    ext = "pdf"
                else:
                    return jsonify({"answer": "Unsupported file type in document generation."})
                download_url = store_file_temporarily(file_stream, ext)
                answer = f"Document generated. You can download it here: {download_url}"
                return jsonify({"answer": answer})

            elif fname == "analyze_data":
                query_analysis = args.get("query")
                if not query_analysis:
                    return jsonify({"answer": "Failed to parse data analysis query."})
                analysis_result = analyze_data(query_analysis)
                answer = f"Data Analysis Result:\n{analysis_result.get('analysis_result', '')}"
                if analysis_result.get("plot_inline"):
                    answer += f"<br>{analysis_result['plot_inline']}"
                return jsonify({"answer": answer})
            else:
                return jsonify({"answer": "Function call not recognized."})
        else:
            answer = message_response.get("content")
            return jsonify({"answer": answer})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/upload", methods=["POST"])
def upload():
    """
    Handles file uploads. Supports PDF, DOCX, XLSX, PPTX, and CSV.
    For Excel files, if a form field 'sheet' is provided, that sheet is used.
    The function extracts text (for retrieval) and, if applicable, reads data into a DataFrame for analysis.

    Returns:
      JSON: The extracted text or an error message.
    """
    uploaded_file = request.files.get("file")
    if not uploaded_file:
        return jsonify({"error": "No file uploaded"}), 400
    filename = uploaded_file.filename
    if filename == "":
        return jsonify({"error": "No selected file"}), 400
    ALLOWED_EXTENSIONS = {"pdf", "docx", "xlsx", "pptx", "csv"}
    extension = filename.rsplit(".", 1)[-1].lower()
    if extension not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"File type .{extension} not allowed"}), 400
    text = ""
    # For Excel files, check for an optional 'sheet' parameter in the form.
    if extension == "pdf":
        text = extract_text_from_pdf(uploaded_file)
    elif extension == "docx":
        text = extract_text_from_docx(uploaded_file)
    elif extension == "xlsx":
        sheet = request.form.get("sheet")  # Optional: user can supply sheet name
        text = extract_text_from_excel(uploaded_file, sheet=sheet)
    elif extension == "csv":
        text = extract_text_from_csv(uploaded_file)
    elif extension == "pptx":
        text = extract_text_from_pptx(uploaded_file)
    else:
        return jsonify({"error": "Unsupported file type"}), 400

    create_vector_store(text)
    return jsonify({"text": text})

@app.route("/download/<filename>")
def download_file(filename):
    """
    Serves a file from temporary storage for download.

    Args:
      filename (str): The name of the file to download.

    Returns:
      The file as an attachment, or a 404 if not found.
    """
    file_path = os.path.join("temp_files", filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        return "File not found", 404

if __name__ == "__main__":
    app.run(debug=True)
